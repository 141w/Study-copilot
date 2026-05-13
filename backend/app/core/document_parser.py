"""
统一文档解析模块 - 支持PDF、DOCX、PPTX格式

架构：
- BaseParser: 解析器抽象基类
- PDFParser: PDF解析实现 (基于 IBM Docling)
- DOCXParser: Word文档解析实现
- PPTXParser: PowerPoint解析实现
- DocumentParser: 统一解析器工厂

数据结构调整 (2025-04-22):
- PDF 解析改用 Docling，输出 Markdown 格式
- 支持复杂表格、多栏排版、扫描件 OCR
- pages 中的 text 字段现为 Markdown 格式

数据结构统一为：
{
    "metadata": {
        "title": "标题",
        "author": "作者", 
        "page_count": int,
        "file_type": "pdf|docx|pptx"
    },
    "pages": [
        {"page": 1, "text": "页面文本内容(Markdown格式)", "images": []},
        ...
    ]
}

作者：AI全栈工程师
"""

from abc import ABC, abstractmethod
from typing import List, Dict, Any, Optional
from pathlib import Path
import os
import asyncio
import logging
from concurrent.futures import ThreadPoolExecutor

from docling.document_converter import DocumentConverter, PdfFormatOption
from docling.datamodel.base_models import InputFormat
from docling.datamodel.pipeline_options import PdfPipelineOptions

logger = logging.getLogger(__name__)

# 线程池，用于 CPU 密集型的 Docling 解析
_executor = ThreadPoolExecutor(max_workers=2)


class BaseParser(ABC):
    """文档解析器抽象基类"""

    @property
    @abstractmethod
    def supported_extensions(self) -> List[str]:
        """支持的文件扩展名"""
        pass

    @abstractmethod
    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析文档"""
        pass

    @abstractmethod
    async def extract_text(self, file_path: str) -> str:
        """提取纯文本"""
        pass

    @abstractmethod
    async def extract_pages(self, file_path: str) -> List[Dict[str, Any]]:
        """提取分页内容"""
        pass

    def is_supported(self, filename: str) -> bool:
        """检查是否支持该文件类型"""
        ext = Path(filename).suffix.lower()
        return ext in self.supported_extensions


class PDFParser(BaseParser):
    """
    PDF文档解析器 - 基于 IBM Docling
    
    特点：
    - 支持复杂表格、多栏排版
    - 支持扫描件 OCR (需要足够内存)
    - 输出 Markdown 格式
    - 使用 asyncio.to_thread 避免阻塞事件循环
    - 智能回退：内存不足时回退到 PyMuPDF
    """

    # 大文档阈值（页数），超过此值不使用 OCR
    LARGE_DOC_THRESHOLD = 30
    
    # 小扫描件阈值，只有小于此页数才尝试 Docling OCR
    OCR_THRESHOLD = 10

    def __init__(self):
        self._converter = None

    @property
    def supported_extensions(self) -> List[str]:
        return [".pdf"]

    def _check_pdf_has_text(self, file_path: str) -> bool:
        """检查 PDF 是否已有嵌入的文本层"""
        try:
            import fitz
            doc = fitz.open(file_path)
            has_text = False
            for page in doc[:3]:  # 检查前几页
                text = page.get_text().strip()
                if text and len(text) > 50:
                    has_text = True
                    break
            doc.close()
            return has_text
        except Exception:
            return False

    def _get_converter(self, enable_ocr: bool = False) -> DocumentConverter:
        """延迟初始化 DocumentConverter"""
        if self._converter is None:
            # 根据是否需要 OCR 配置
            pipeline_options = PdfPipelineOptions()
            pipeline_options.do_ocr = enable_ocr
            pipeline_options.do_table_structure = enable_ocr
            
            self._converter = DocumentConverter(
                format_options={
                    InputFormat.PDF: PdfFormatOption(
                        pipeline_options=pipeline_options
                    ),
                }
            )
        return self._converter

    def _convert_with_docling(self, file_path: str, enable_ocr: bool = False) -> Dict[str, Any]:
        """使用 Docling 转换"""
        try:
            import fitz
            # 获取页数
            temp_doc = fitz.open(file_path)
            page_count = len(temp_doc)
            temp_doc.close()
            
            converter = self._get_converter(enable_ocr)
            result = converter.convert(file_path)
            
            markdown_text = result.document.export_to_markdown()
            
            doc = result.document
            metadata = {
                "title": getattr(doc, 'title', '') or '',
                "author": getattr(doc, 'authors', '') or '',
                "subject": getattr(doc, 'subject', '') or '',
                "creator": getattr(doc, 'creator', '') or '',
                "page_count": page_count,
                "file_type": "pdf"
            }
            
            logger.info(f"Docling succeeded: {len(markdown_text)} chars")
            
            return {
                "metadata": metadata,
                "markdown": markdown_text,
                "success": True
            }
            
        except Exception as e:
            logger.error(f"Docling 解析失败: {file_path}, error: {e}")
            return {
                "metadata": {"file_type": "pdf"},
                "markdown": "",
                "success": False,
                "error": str(e)
            }

    def _convert_fallback(self, file_path: str) -> Dict[str, Any]:
        """回退方案：使用 PyMuPDF"""
        try:
            import fitz
            doc = fitz.open(file_path)
            
            page_count = len(doc)
            text_content = []
            for page_num, page in enumerate(doc):
                text = page.get_text("text").strip()
                if text:
                    text_content.append(text)
            
            doc.close()
            
            markdown_text = "\n\n".join(text_content)
            
            return {
                "metadata": {
                    "title": "",
                    "author": "",
                    "subject": "",
                    "creator": "",
                    "page_count": page_count,
                    "file_type": "pdf",
                    "_fallback": True
                },
                "markdown": markdown_text,
                "success": True
            }
        except Exception as e:
            logger.error(f"PyMuPDF 回退也失败: {file_path}, error: {e}")
            return {
                "metadata": {"file_type": "pdf"},
                "markdown": "",
                "success": False,
                "error": str(e)
            }

    def _split_markdown_pages(self, markdown_text: str, page_count: int) -> List[Dict[str, Any]]:
        """
        将 Markdown 文本分割成页面块
        
        注意: Docling 返回的是完整的 Markdown，不是按页分割的。
        为了兼容下游 chunker，我们有几种策略：
        1. 整个文档作为一页 - 简单但下游处理慢
        2. 按 Markdown 标题层级分割 - 保持语义
        3. 按固定段落数分页 - 均衡
        
        这里采用策略1：将整个 Markdown 作为一页返回，
        让下游的 chunker 处理具体的分块
        """
        if not markdown_text or not markdown_text.strip():
            return [{
                "page": 1,
                "text": "",
                "images": []
            }]
        
        # 方法1: 直接返回整个文档（推荐，让 chunker 处理）
        # 这样可以保持 Markdown 结构的完整性
        return [{
            "page": 1,
            "text": markdown_text,
            "images": []
        }]
        
        # 方法2: 按标题拆分 (可选，如果文档很大可以启用)
        # import re
        # if len(markdown_text) > 50000:  # 超过50KB时按标题拆分
        #     # 按 ## 标题拆分
        #     sections = re.split(r'^##\s+', markdown_text, flags=re.MULTILINE)
        #     pages = []
        #     for idx, section in enumerate(sections):
        #         if section.strip():
        #             pages.append({
        #                 "page": idx + 1,
        #                 "text": ("## " + section).strip() if idx > 0 else section.strip(),
        #                 "images": []
        #             })
        #     if pages:
        #         return pages
        
        # return [{
        #     "page": 1,
        #     "text": markdown_text,
        #     "images": []
        # }]

    async def parse(self, file_path: str) -> Dict[str, Any]:
        """
        解析 PDF 文件
        
        智能策略：
        1. 先检查 PDF 是否有嵌入文本
        2. 有文本 → 正常解析（PyMuPDF 或 Docling）
        3. 无文本（扫描件）：
           - 小文档(<10页) → 尝试 Docling OCR
           - 大文档(>=10页) → 提示需要可搜索 PDF
        4. 使用 asyncio.to_thread 避免阻塞事件循环
        """
        loop = asyncio.get_event_loop()
        page_count = 1
        
        try:
            # 先检查页数和是否有文本
            import fitz
            temp_doc = fitz.open(file_path)
            page_count = len(temp_doc)
            
            # 检查前几页是否有文本
            has_text = False
            for page in temp_doc[:min(3, page_count)]:
                if page.get_text().strip():
                    has_text = True
                    break
            temp_doc.close()
            
            logger.info(f"PDF: {page_count} pages, has_text={has_text}")
            
            if has_text:
                # 有文本层的 PDF
                if page_count > self.LARGE_DOC_THRESHOLD:
                    # 大文档直接用 PyMuPDF
                    logger.info(f"Large doc with text, using PyMuPDF")
                    result = await loop.run_in_executor(
                        _executor,
                        self._convert_fallback,
                        file_path
                    )
                else:
                    # 小文档用 Docling 获得更好解析
                    logger.info(f"Small doc with text, using Docling")
                    result = await loop.run_in_executor(
                        _executor,
                        self._convert_with_docling,
                        file_path,
                        False
                    )
                    if not result.get("success", False):
                        result = await loop.run_in_executor(
                            _executor,
                            self._convert_fallback,
                            file_path
                        )
            else:
                # 扫描件 PDF（无文本层）
                if page_count < self.OCR_THRESHOLD:
                    # 小扫描件，尝试 Docling OCR
                    logger.info(f"Scanned PDF (<{self.OCR_THRESHOLD}p), trying Docling OCR")
                    result = await loop.run_in_executor(
                        _executor,
                        self._convert_with_docling,
                        file_path,
                        True  # enable_ocr=True
                    )
                    if not result.get("success", False):
                        result = await loop.run_in_executor(
                            _executor,
                            self._convert_fallback,
                            file_path
                        )
                else:
                    # 大扫描件，必须先转换
                    raise ValueError(
                        f"检测到扫描件 PDF ({page_count} 页)，无法直接提取文字。"
                        "请先将扫描件转换为可搜索 PDF（包含文本层），"
                        "或使用支持 OCR 的工具预处理后再上传。"
                    )
            
            if not result.get("success", False):
                error_msg = result.get("error", "未知错误")
                raise ValueError(f"文档解析失��: {error_msg}")
            
            # 转换为 pages 格式
            markdown_text = result.get("markdown", "")
            
            # 再次检查是否提取到文本
            if not markdown_text or len(markdown_text.strip()) < 50:
                raise ValueError(
                    "无法提取有效文本内容。"
                    "如果这是扫描件，请先转换为可搜索 PDF。"
                )
            
            pages = self._split_markdown_pages(markdown_text, page_count)
            
            logger.info(f"PDF parsed: {len(pages)} pages, {len(markdown_text)} chars")
            
            return {
                "metadata": result["metadata"],
                "pages": pages
            }
        except ValueError:
            raise
        except Exception as e:
            logger.error(f"PDF 解析异常: {file_path}, error: {e}")
            raise ValueError(f"文档解析失败: {str(e)}")

    async def extract_text(self, file_path: str) -> str:
        """提取纯文本 (Markdown 格式)"""
        result = await self.parse(file_path)
        return "\n\n".join([page["text"] for page in result["pages"]])

    async def extract_pages(self, file_path: str) -> List[Dict[str, Any]]:
        """提取分页内容"""
        result = await self.parse(file_path)
        return result["pages"]


class DOCXParser(BaseParser):
    """Word文档(.docx)解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".docx", ".doc"]

    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析Word文档"""
        from docx import Document
        
        doc = Document(file_path)
        
        # 提取元数据
        core_props = doc.core_properties
        metadata = {
            "title": core_props.title or "",
            "author": core_props.author or "",
            "subject": core_props.subject or "",
            "creator": core_props.author or "",
            "page_count": len(doc.paragraphs),  # Word没有页码概念，用段落数近似
            "file_type": "docx"
        }

        # 按段落提取文本
        text_content = []
        current_page_text = []
        page_num = 1
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if text:
                # 检查是否是新段落的开始（可作为新页的标志）
                # 这里简单按段落数分组，每20个段落作为一页
                current_page_text.append(text)
                
                if len(current_page_text) >= 20:
                    text_content.append({
                        "page": page_num,
                        "text": "\n".join(current_page_text),
                        "images": []
                    })
                    page_num += 1
                    current_page_text = []

        # 处理剩余段落
        if current_page_text:
            text_content.append({
                "page": page_num,
                "text": "\n".join(current_page_text),
                "images": []
            })

        # 也处理表格中的文本
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        table_text.append(cell.text.strip())
            
            if table_text:
                text_content.append({
                    "page": page_num,
                    "text": "\n".join(table_text),
                    "images": []
                })

        return {"metadata": metadata, "pages": text_content}

    async def extract_text(self, file_path: str) -> str:
        """提取纯文本"""
        result = await self.parse(file_path)
        return "\n\n".join([page["text"] for page in result["pages"]])

    async def extract_pages(self, file_path: str) -> List[Dict[str, Any]]:
        """提取分页内容"""
        result = await self.parse(file_path)
        return result["pages"]


class PPTXParser(BaseParser):
    """PowerPoint(.pptx)解析器"""

    @property
    def supported_extensions(self) -> List[str]:
        return [".pptx", ".ppt"]

    async def parse(self, file_path: str) -> Dict[str, Any]:
        """解析PowerPoint文件"""
        from pptx import Presentation
        
        prs = Presentation(file_path)
        
        metadata = {
            "title": prs.core_properties.title or "",
            "author": prs.core_properties.author or "",
            "subject": prs.core_properties.subject or "",
            "creator": prs.core_properties.author or "",
            "page_count": len(prs.slides),
            "file_type": "pptx"
        }

        # 按幻灯片提取文本
        text_content = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_text = []
            
            # 从形状中提取文本
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
                
                # 从表格中提取
                if shape.has_table:
                    for row in shape.table.rows:
                        for cell in row.cells:
                            if cell.text_frame.text.strip():
                                slide_text.append(cell.text_frame.text.strip())
            
            if slide_text:
                text_content.append({
                    "page": slide_num + 1,
                    "text": "\n".join(slide_text),
                    "images": []
                })
            else:
                # 空幻灯片也保留
                text_content.append({
                    "page": slide_num + 1,
                    "text": "",
                    "images": []
                })

        return {"metadata": metadata, "pages": text_content}

    async def extract_text(self, file_path: str) -> str:
        """提取纯文本"""
        result = await self.parse(file_path)
        texts = []
        for page in result["pages"]:
            if page["text"].strip():
                texts.append(f"[幻灯片 {page['page']}]\n{page['text']}")
        return "\n\n".join(texts)

    async def extract_pages(self, file_path: str) -> List[Dict[str, Any]]:
        """提取分页内容"""
        result = await self.parse(file_path)
        return result["pages"]


class DocumentParser:
    """
    统一文档解析器工厂
    
    根据文件类型自动选择合适的解析器
    """

    def __init__(self):
        self._parsers = {
            ".pdf": PDFParser(),
            ".docx": DOCXParser(),
            ".doc": DOCXParser(),
            ".pptx": PPTXParser(),
            ".ppt": PPTXParser(),
        }

    @property
    def supported_extensions(self) -> List[str]:
        """所有支持的文件扩展名"""
        exts = []
        for parser in self._parsers.values():
            exts.extend(parser.supported_extensions)
        return list(set(exts))

    def is_supported(self, filename: str) -> bool:
        """检查是否支持"""
        ext = Path(filename).suffix.lower()
        return ext in self.supported_extensions

    def get_parser(self, file_path: str) -> Optional[BaseParser]:
        """获取合适的解析器"""
        ext = Path(file_path).suffix.lower()
        return self._parsers.get(ext)

    async def parse(self, file_path: str) -> Dict[str, Any]:
        """统一解析接口"""
        parser = self.get_parser(file_path)
        if parser is None:
            raise ValueError(f"不支持的文件类型: {Path(file_path).suffix}")
        return await parser.parse(file_path)

    async def extract_text(self, file_path: str) -> str:
        """提取纯文本"""
        parser = self.get_parser(file_path)
        if parser is None:
            raise ValueError(f"不支持的文件类型: {Path(file_path).suffix}")
        return await parser.extract_text(file_path)

    async def extract_pages(self, file_path: str) -> List[Dict[str, Any]]:
        """提取分页内容"""
        parser = self.get_parser(file_path)
        if parser is None:
            raise ValueError(f"不支持的文件类型: {Path(file_path).suffix}")
        return await parser.extract_pages(file_path)


# 导出默认实例
document_parser = DocumentParser()